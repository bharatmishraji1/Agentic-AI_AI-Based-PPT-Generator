import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


class SlideGenerator:
    def __init__(self, output_dir="output/"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

        self.prs = Presentation()

    #  Title Slide
    def add_title_slide(self, title, subtitle=""):
        layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)

        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    #  Content Slide
    def add_content_slide(self, title, bullets):
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)

        slide.shapes.title.text = title

        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for bullet in bullets:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0, 0, 0)

    # Conclusion Slide
    def add_conclusion_slide(self, title, bullets):
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)

        slide.shapes.title.text = title

        tf = slide.placeholders[1].text_frame
        tf.clear()

        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(26)
            p.font.bold = True

    #  Main Generator
    def generate_presentation(self, slides_data, filename="generated_presentation.pptx"):
        for i, slide in enumerate(slides_data):
            slide_type = slide.get("type", "content")

            if slide_type == "intro":
                self.add_title_slide(slide["title"], "Auto Generated Presentation")

            elif slide_type == "conclusion":
                self.add_conclusion_slide(slide["title"], slide["content"])

            else:
                self.add_content_slide(slide["title"], slide["content"])

        output_path = os.path.join(self.output_dir, filename)
        self.prs.save(output_path)

        print(f"\n Presentation saved at: {output_path}")


# Example usage
if __name__ == "__main__":

    #  This is what your AI should generate
    slides_data = [
        {
            "title": "Artificial Intelligence",
            "type": "intro",
            "content": []
        },
        {
            "title": "What is AI?",
            "type": "content",
            "content": [
                "AI simulates human intelligence",
                "Machines learn from data",
                "Used across multiple industries"
            ]
        },
        {
            "title": "Applications of AI",
            "type": "content",
            "content": [
                "Healthcare (disease prediction)",
                "Finance (fraud detection)",
                "Marketing (customer targeting)",
                "Self-driving vehicles"
            ]
        },
        {
            "title": "Benefits",
            "type": "content",
            "content": [
                "Automation of tasks",
                "Improved efficiency",
                "Better decision making"
            ]
        },
        {
            "title": "Conclusion",
            "type": "conclusion",
            "content": [
                "AI is transforming industries",
                "Future is driven by intelligent systems",
                "Adoption will continue to grow"
            ]
        }
    ]

    generator = SlideGenerator()
    generator.generate_presentation(slides_data)